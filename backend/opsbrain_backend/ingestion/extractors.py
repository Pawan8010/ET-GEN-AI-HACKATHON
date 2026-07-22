import csv
import io
import zipfile
from pathlib import Path

from bs4 import BeautifulSoup
from docx import Document as WordDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader


MAX_ARCHIVE_FILES = 100
MAX_ARCHIVE_BYTES = 50 * 1024 * 1024


def extract_text(filename: str, content: bytes) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix in {".txt", ".md", ".log"}:
        return content.decode("utf-8", errors="replace")
    if suffix == ".csv":
        rows = csv.reader(io.StringIO(content.decode("utf-8-sig", errors="replace")))
        return "\n".join(" | ".join(cell.strip() for cell in row) for row in rows)
    if suffix in {".html", ".htm"}:
        return BeautifulSoup(content, "html.parser").get_text("\n", strip=True)
    if suffix == ".pdf":
        return "\n\n".join(page.extract_text() or "" for page in PdfReader(io.BytesIO(content)).pages)
    if suffix == ".docx":
        document = WordDocument(io.BytesIO(content))
        paragraphs = [paragraph.text for paragraph in document.paragraphs]
        tables = [" | ".join(cell.text for cell in row.cells) for table in document.tables for row in table.rows]
        return "\n".join(paragraphs + tables)
    if suffix == ".pptx":
        presentation = Presentation(io.BytesIO(content))
        return "\n".join(
            shape.text for slide in presentation.slides for shape in slide.shapes
            if hasattr(shape, "text") and shape.text
        )
    if suffix in {".xlsx", ".xlsm"}:
        workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        return "\n".join(
            " | ".join("" if value is None else str(value) for value in row)
            for sheet in workbook.worksheets for row in sheet.iter_rows(values_only=True)
        )
    if suffix == ".zip":
        return _extract_zip(content)
    raise ValueError(f"Unsupported document format: {suffix or 'unknown'}")


def _extract_zip(content: bytes) -> str:
    output: list[str] = []
    total = 0
    with zipfile.ZipFile(io.BytesIO(content)) as archive:
        members = archive.infolist()
        if len(members) > MAX_ARCHIVE_FILES:
            raise ValueError("Archive contains too many files")
        for member in members:
            path = Path(member.filename)
            if member.is_dir():
                continue
            if path.is_absolute() or ".." in path.parts:
                raise ValueError("Archive contains an unsafe path")
            total += member.file_size
            if total > MAX_ARCHIVE_BYTES:
                raise ValueError("Archive expanded size exceeds the limit")
            try:
                extracted = extract_text(path.name, archive.read(member))
            except ValueError:
                continue
            output.append(f"# {member.filename}\n{extracted}")
    return "\n\n".join(output)
